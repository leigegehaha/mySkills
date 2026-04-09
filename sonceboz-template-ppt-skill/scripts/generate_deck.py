from __future__ import annotations

import argparse
import json
import re
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SKILL_DIR = Path(__file__).resolve().parents[1]
TEMPLATE_PATH = SKILL_DIR / "assets" / "template.pptx"
SVG_FACTORY_PATH = SKILL_DIR / "scripts" / "svg_factory.py"
WEB_IMAGE_SEARCH_PATH = SKILL_DIR.parent / "web-image-research" / "scripts" / "search_images.py"

RED = RGBColor(192, 0, 0)
RED_DARK = RGBColor(145, 0, 0)
DARK = RGBColor(51, 51, 51)
MID = RGBColor(102, 102, 102)
LIGHT = RGBColor(245, 245, 245)
BORDER = RGBColor(217, 217, 217)
WHITE = RGBColor(255, 255, 255)


def remove_all_slides(presentation: Presentation) -> None:
    slide_ids = presentation.slides._sldIdLst
    for slide_id in list(slide_ids):
        presentation.part.drop_rel(slide_id.rId)
        slide_ids.remove(slide_id)


def clear_text_frame(text_frame, margin: float = 0.08, vertical_anchor=MSO_ANCHOR.TOP) -> None:
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(margin)
    text_frame.margin_right = Inches(margin)
    text_frame.margin_top = Inches(margin)
    text_frame.margin_bottom = Inches(margin)
    text_frame.vertical_anchor = vertical_anchor


def style_paragraph(paragraph, text: str, size: int = 16, color: RGBColor = DARK, bold: bool = False, align=PP_ALIGN.LEFT, space_after: int = 8) -> None:
    paragraph.text = text
    paragraph.alignment = align
    paragraph.space_after = Pt(space_after)
    paragraph.line_spacing = 1.05
    run = paragraph.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_textbox(slide, left: float, top: float, width: float, height: float, paragraphs: list[dict], margin: float = 0.08, vertical_anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    clear_text_frame(tf, margin=margin, vertical_anchor=vertical_anchor)
    for idx, spec in enumerate(paragraphs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        style_paragraph(
            p,
            spec["text"],
            size=spec.get("size", 16),
            color=spec.get("color", DARK),
            bold=spec.get("bold", False),
            align=spec.get("align", PP_ALIGN.LEFT),
            space_after=spec.get("space_after", 8),
        )
    return box


def add_card(slide, left: float, top: float, width: float, height: float, fill: RGBColor = WHITE, line: RGBColor = BORDER):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_picture_contain(slide, image_path: Path, left: float, top: float, width: float, height: float):
    add_card(slide, left, top, width, height, fill=LIGHT, line=BORDER)
    with Image.open(image_path) as image:
        ratio = image.width / image.height
    box_ratio = width / height
    if ratio >= box_ratio:
        pic_width = Inches(width)
        pic_height = int(pic_width / ratio)
    else:
        pic_height = Inches(height)
        pic_width = int(pic_height * ratio)
    slide.shapes.add_picture(
        str(image_path),
        Inches(left) + (Inches(width) - pic_width) // 2,
        Inches(top) + (Inches(height) - pic_height) // 2,
        width=pic_width,
        height=pic_height,
    )


def generate_visual_assets(source: dict, purpose: str, output_path: Path) -> list[Path]:
    if not SVG_FACTORY_PATH.exists():
        return []
    namespace = {}
    exec(SVG_FACTORY_PATH.read_text(encoding="utf-8"), namespace)
    chooser = namespace.get("choose_visuals")
    if not chooser:
        return []
    payload = dict(source)
    payload["output_dir"] = str((output_path.parent / f"{output_path.stem}_visuals").resolve())
    payload["title"] = payload.get("title") or output_path.stem
    pairs = chooser(payload, purpose)
    return [Path(png) for _, png in pairs if Path(png).exists()]


def generate_external_image_records(source: dict, purpose: str, output_path: Path, richness: int) -> list[dict]:
    if richness <= 0 or not WEB_IMAGE_SEARCH_PATH.exists():
        return []
    namespace = {"__file__": str(WEB_IMAGE_SEARCH_PATH)}
    exec(WEB_IMAGE_SEARCH_PATH.read_text(encoding="utf-8"), namespace)
    search_and_collect = namespace.get("search_and_collect")
    if not search_and_collect:
        return []
    out_dir = (output_path.parent / f"{output_path.stem}_web_images").resolve()
    manifest = search_and_collect(source, purpose, out_dir, richness)
    return [item for item in manifest.get("images", []) if item.get("local_path") and Path(item["local_path"]).exists()]


def infer_local_image_role(path: Path, purpose: str, kind: str) -> str:
    name = path.stem.lower().replace("-", " ").replace("_", " ")
    if kind == "generated":
        if "timeline" in name:
            return "method"
        if any(token in name for token in ["metric", "chart", "grid", "result", "trend"]):
            return "results"
        return "visual"
    if purpose == "test-report":
        if any(token in name for token in ["test", "bench", "setup", "lab", "thermal", "camera", "fixture", "measure"]):
            return "evidence"
        if any(token in name for token in ["hand", "gripper", "robot", "product", "device", "finger"]):
            return "product"
        return "source"
    if purpose == "market-research":
        if any(token in name for token in ["product", "robot", "device", "component"]):
            return "product"
        return "context"
    if any(token in name for token in ["meeting", "workshop", "team", "office", "discussion"]):
        return "meeting"
    return "context"


def build_image_catalog(source_images: list[Path], external_images: list[dict], generated_visuals: list[Path], purpose: str) -> list[dict]:
    catalog = []
    for idx, path in enumerate(source_images):
        if not path.exists():
            continue
        catalog.append(
            {
                "path": path,
                "kind": "source",
                "role": infer_local_image_role(path, purpose, "source"),
                "title": path.stem,
                "score": 120 - idx * 3,
            }
        )
    for idx, item in enumerate(external_images):
        path = Path(item["local_path"])
        if not path.exists():
            continue
        item_kind = "ai" if item.get("kind") == "ai" else "external"
        catalog.append(
            {
                "path": path,
                "kind": item_kind,
                "role": item.get("role") or "context",
                "title": item.get("title", path.stem),
                "score": 80 + float(item.get("score", 0)) - idx * 0.2,
                "source_name": item.get("source", "web"),
            }
        )
    for idx, path in enumerate(generated_visuals):
        if not path.exists():
            continue
        role = infer_local_image_role(path, purpose, "generated")
        base = 95 if role in {"results", "method"} else 68
        catalog.append(
            {
                "path": path,
                "kind": "generated",
                "role": role,
                "title": path.stem,
                "score": base - idx * 2,
            }
        )
    return catalog


def pick_catalog_image(
    catalog: list[dict],
    used: set[Path],
    preferred_roles: tuple[str, ...] = (),
    preferred_kinds: tuple[str, ...] = (),
    avoid_roles: tuple[str, ...] = (),
    avoid_kinds: tuple[str, ...] = (),
):
    ranked = []
    for item in catalog:
        path = item["path"]
        if path in used:
            continue
        score = float(item.get("score", 0))
        role = item.get("role", "")
        kind = item.get("kind", "")
        if preferred_roles:
            score += 28 - preferred_roles.index(role) * 4 if role in preferred_roles else -8
        if preferred_kinds:
            score += 14 - preferred_kinds.index(kind) * 2 if kind in preferred_kinds else -4
        if role in avoid_roles:
            score -= 24
        if kind in avoid_kinds:
            score -= 14
        ranked.append((score, item))
    if not ranked:
        return None
    ranked.sort(key=lambda pair: pair[0], reverse=True)
    chosen = ranked[0][1]
    used.add(chosen["path"])
    return chosen


def image_caption(record: dict | None) -> str:
    if not record:
        return ""
    return {
        "product": "Representative product visual",
        "evidence": "Representative test setup visual",
        "method": "Generated method-support visual",
        "results": "Generated results-support visual",
        "meeting": "Meeting context visual",
        "context": "Supporting context visual",
        "source": "Representative source image",
        "visual": "Supporting visual",
    }.get(record.get("role", ""), "Supporting visual")


def choose_test_report_images(source_images: list[Path], external_images: list[dict], generated_visuals: list[Path]) -> dict[str, dict | None]:
    catalog = build_image_catalog(source_images, external_images, generated_visuals, "test-report")
    used: set[Path] = set()
    overview = pick_catalog_image(
        catalog,
        used,
        preferred_roles=("product", "source", "context", "evidence"),
        preferred_kinds=("source", "external", "generated"),
        avoid_roles=("results", "method"),
    )
    method = pick_catalog_image(
        catalog,
        used,
        preferred_roles=("evidence", "method", "product", "source"),
        preferred_kinds=("source", "external", "generated"),
        avoid_roles=("results",),
    )
    results = pick_catalog_image(
        catalog,
        used,
        preferred_roles=("results", "method", "product", "evidence"),
        preferred_kinds=("generated", "external", "source"),
    )
    return {"overview": overview, "method": method, "results": results}


def choose_market_images(source_images: list[Path], external_images: list[dict], generated_visuals: list[Path]) -> dict[str, dict | None]:
    catalog = build_image_catalog(source_images, external_images, generated_visuals, "market-research")
    used: set[Path] = set()
    landscape = pick_catalog_image(catalog, used, preferred_roles=("context", "product"), preferred_kinds=("external", "source", "generated"))
    implication = pick_catalog_image(catalog, used, preferred_roles=("product", "context"), preferred_kinds=("external", "source", "generated"))
    return {"landscape": landscape, "implication": implication}


def choose_meeting_images(source_images: list[Path], external_images: list[dict], generated_visuals: list[Path]) -> dict[str, dict | None]:
    catalog = build_image_catalog(source_images, external_images, generated_visuals, "meeting-summary")
    used: set[Path] = set()
    discussion = pick_catalog_image(catalog, used, preferred_roles=("meeting", "context"), preferred_kinds=("external", "source", "generated"))
    return {"discussion": discussion}


def remaining_catalog_paths(source_images: list[Path], external_images: list[dict], generated_visuals: list[Path], purpose: str, used_records: list[dict | None]) -> list[Path]:
    used_paths = {record["path"] for record in used_records if record}
    catalog = build_image_catalog(source_images, external_images, generated_visuals, purpose)
    ranked = sorted(catalog, key=lambda item: float(item.get("score", 0)), reverse=True)
    return [item["path"] for item in ranked if item["path"] not in used_paths]


def choose_image(primary: list[Path], external: list[Path], generated: list[Path], index: int) -> Path | None:
    for group in (primary, external, generated):
        if index < len(group):
            return group[index]
    return None


def add_image_mosaic(slide, images: list[Path], left: float, top: float, width: float, height: float, max_items: int = 3) -> None:
    items = images[:max_items]
    if not items:
        return
    gap = 0.12
    cell_width = (width - gap * (len(items) - 1)) / len(items)
    for idx, image_path in enumerate(items):
        add_picture_contain(slide, image_path, left + idx * (cell_width + gap), top, cell_width, height)


def add_title(slide, title: str) -> None:
    placeholder = slide.placeholders[0]
    placeholder.text = title
    tf = placeholder.text_frame
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        for r in p.runs:
            r.font.name = "Arial"
            r.font.size = Pt(28)
            r.font.bold = True
            r.font.color.rgb = RED


def add_kicker(slide, text: str) -> None:
    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(1.03), Inches(2.0), Inches(0.32))
    tag.fill.solid()
    tag.fill.fore_color.rgb = RED
    tag.line.color.rgb = RED
    tf = tag.text_frame
    clear_text_frame(tf, margin=0.02, vertical_anchor=MSO_ANCHOR.MIDDLE)
    style_paragraph(tf.paragraphs[0], text, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER, space_after=0)


def add_placeholder(slide, left: float, top: float, width: float, height: float, heading: str, note: str, red: bool = True) -> None:
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(252, 252, 252)
    box.line.color.rgb = RED if red else BORDER
    box.line.width = Pt(1.8 if red else 1.2)
    box.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    tf = box.text_frame
    clear_text_frame(tf, margin=0.1, vertical_anchor=MSO_ANCHOR.MIDDLE)
    style_paragraph(tf.paragraphs[0], heading, size=18 if red else 15, color=RED if red else DARK, bold=True, align=PP_ALIGN.CENTER, space_after=4)
    p = tf.add_paragraph()
    style_paragraph(p, note, size=11, color=MID, align=PP_ALIGN.CENTER, space_after=0)


def polish_title(raw: str, purpose: str) -> str:
    if raw.strip():
        return raw.strip()
    defaults = {
        "test-report": "Validation Report",
        "market-research": "Market Research Summary",
        "meeting-summary": "Meeting Summary",
    }
    return defaults[purpose]


def polish_sentence(line: str) -> str:
    text = re.sub(r"\s+", " ", line).strip(" -")
    if not text:
        return ""
    replacements = {
        "vs ": "versus ",
        "max ": "maximum ",
        "min ": "minimum ",
        "temp ": "temperature ",
        "grip ": "grasp ",
    }
    lower = text.lower()
    for old, new in replacements.items():
        lower = lower.replace(old, new)
    text = lower[:1].upper() + lower[1:]
    if text[-1] not in ".!?":
        text += "."
    return text


def is_meta_line(line: str) -> bool:
    lower = line.strip().lower()
    instruction_prefixes = (
        "purpose:",
        "title:",
        "this deck should",
        "this document should",
        "please keep the tone",
        "please preserve",
        "please polish",
        "add rich visual support",
        "reserve placeholders",
    )
    return lower.startswith(instruction_prefixes) or lower in {"key source facts:", "key source facts"}


def dedupe_lines(lines: list[str]) -> list[str]:
    seen: set[str] = set()
    output: list[str] = []
    for line in lines:
        norm = line.strip().lower()
        if not norm or norm in seen:
            continue
        seen.add(norm)
        output.append(line.strip())
    return output


def summarize_lines(lines: list[str], limit: int = 8) -> list[str]:
    return dedupe_lines([polish_sentence(line) for line in lines if len(line.strip()) > 3 and not is_meta_line(line)])[:limit]


def choose_lines_for_purpose(text: str, purpose: str) -> dict[str, list[str]]:
    lines = [line.strip() for line in text.splitlines() if line.strip() and not is_meta_line(line)]
    if purpose == "test-report":
        summary = summarize_lines(lines, 8)
        metrics = [line for line in lines if re.search(r"\d", line)][:10]
        methods = [line for line in lines if any(token in line.lower() for token in ["method", "setup", "test", "instrument", "measure", "protocol"])]
        return {
            "summary": summary[:4],
            "methods": summarize_lines(methods or lines, 4),
            "results": summarize_lines(metrics + lines, 6),
            "conclusion": summarize_lines(lines[::-1], 4),
        }
    if purpose == "market-research":
        insights = summarize_lines(lines, 10)
        return {
            "summary": insights[:4],
            "landscape": insights[1:5],
            "findings": insights[4:8],
            "recommendation": insights[6:10],
        }
    insights = summarize_lines(lines, 10)
    return {
        "summary": insights[:4],
        "discussion": insights[2:6],
        "decisions": insights[4:8],
        "next_steps": insights[6:10],
    }


def filtered_metric_facts(facts: list[dict]) -> list[dict]:
    filtered = []
    for fact in facts:
        label = fact["label"].strip()
        lower = label.lower()
        unit = fact.get("unit", "").strip()
        if lower.startswith(("purpose", "title")):
            continue
        if len(label) < 4:
            continue
        if not unit:
            continue
        filtered.append(fact)
    return filtered


def metric_cards_from_facts(slide, facts: list[dict], start_left: float = 0.45, top: float = 5.62, width: float = 2.9, height: float = 0.55, max_cards: int = 4) -> None:
    filtered = filtered_metric_facts(facts)
    for idx, fact in enumerate(filtered[:max_cards]):
        left = start_left + idx * (width + 0.12)
        add_card(slide, left, top, width, height)
        tf = slide.shapes[-1].text_frame
        clear_text_frame(tf, margin=0.06, vertical_anchor=MSO_ANCHOR.MIDDLE)
        style_paragraph(tf.paragraphs[0], f"{fact['value']} {fact['unit']}".strip(), size=19, color=RED, bold=True, align=PP_ALIGN.CENTER, space_after=2)
        p = tf.add_paragraph()
        style_paragraph(p, fact["label"][:42], size=10, color=MID, align=PP_ALIGN.CENTER, space_after=0)


def metric_cards_vertical(slide, facts: list[dict], left: float, top: float, width: float, height: float, max_cards: int = 4) -> None:
    filtered = filtered_metric_facts(facts)[:max_cards]
    if not filtered:
        add_placeholder(slide, left, top, width, height, "METRIC PLACEHOLDER", "Insert measured values or KPI callouts here.", red=False)
        return
    gap = 0.1
    card_height = (height - gap * (len(filtered) - 1)) / len(filtered)
    for idx, fact in enumerate(filtered):
        card_top = top + idx * (card_height + gap)
        add_card(slide, left, card_top, width, card_height)
        tf = slide.shapes[-1].text_frame
        clear_text_frame(tf, margin=0.05, vertical_anchor=MSO_ANCHOR.MIDDLE)
        style_paragraph(tf.paragraphs[0], f"{fact['value']} {fact['unit']}".strip(), size=18, color=RED, bold=True, align=PP_ALIGN.CENTER, space_after=1)
        p = tf.add_paragraph()
        style_paragraph(p, fact["label"][:32], size=9, color=MID, align=PP_ALIGN.CENTER, space_after=0)


def build_cover(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.placeholders[14].text = title.upper()
    slide.placeholders[15].text = subtitle
    for idx in (14, 15):
        tf = slide.placeholders[idx].text_frame
        for p in tf.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = "Arial"
                r.font.color.rgb = WHITE
                if idx == 14:
                    r.font.size = Pt(36)
                    r.font.bold = True
                else:
                    r.font.size = Pt(18)


def build_summary_slide(prs: Presentation, title: str, lines: list[str], kicker: str = "SUMMARY") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title(slide, title)
    add_kicker(slide, kicker)
    add_card(slide, 0.45, 1.55, 12.0, 4.35)
    add_textbox(slide, 0.75, 1.9, 11.3, 3.6, [{"text": f"• {line}", "size": 15} for line in lines[:6]])


def build_thank_you_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[7])
    tf = slide.placeholders[13].text_frame
    clear_text_frame(tf, margin=0.08, vertical_anchor=MSO_ANCHOR.MIDDLE)
    style_paragraph(tf.paragraphs[0], "Thank you", size=30, color=WHITE, bold=True, align=PP_ALIGN.LEFT, space_after=0)


def build_test_report(prs: Presentation, title: str, source: dict, text_blocks: dict[str, list[str]], image_paths: list[Path], external_images: list[dict], output_path: Path, generated_visuals: list[Path], richness: int) -> None:
    image_plan = choose_test_report_images(image_paths, external_images, generated_visuals)
    build_cover(prs, title, "Generated from the bundled Sonceboz template")

    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title(slide, "Executive Summary")
    add_kicker(slide, "REPORT OVERVIEW")
    add_card(slide, 0.45, 1.75, 5.85, 3.75)
    add_card(slide, 6.45, 1.75, 5.95, 3.75)
    add_textbox(slide, 0.7, 1.96, 5.2, 0.3, [{"text": "Assessment Scope", "size": 16, "bold": True, "space_after": 0}], margin=0)
    add_textbox(slide, 0.7, 2.28, 5.2, 2.9, [{"text": f"• {line}", "size": 14} for line in text_blocks["summary"][:4]])
    add_textbox(slide, 6.7, 1.96, 5.2, 0.3, [{"text": "Key Findings", "size": 16, "bold": True, "space_after": 0}], margin=0)
    add_textbox(slide, 6.7, 2.28, 5.1, 2.9, [{"text": f"• {line}", "size": 14} for line in text_blocks["results"][:4]])
    metric_cards_from_facts(slide, source.get("numeric_facts", []))

    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title(slide, "Source Overview")
    add_kicker(slide, "INPUT SNAPSHOT")
    hero_image = image_plan["overview"]
    if hero_image:
        add_picture_contain(slide, hero_image["path"], 0.45, 1.55, 5.3, 3.95)
        add_textbox(slide, 0.6, 5.1, 5.0, 0.24, [{"text": image_caption(hero_image), "size": 10, "color": MID, "align": PP_ALIGN.CENTER, "space_after": 0}], margin=0)
    else:
        add_placeholder(slide, 0.45, 1.55, 5.3, 3.95, "IMAGE PLACEHOLDER", "No suitable source image was extracted. Replace with a product or setup visual.", red=False)
    add_card(slide, 6.0, 1.55, 6.4, 3.95)
    add_textbox(slide, 6.25, 1.82, 5.9, 0.32, [{"text": "Extracted Highlights", "size": 16, "bold": True, "space_after": 0}], margin=0)
    add_textbox(slide, 6.25, 2.18, 5.7, 2.9, [{"text": f"• {line}", "size": 14} for line in (text_blocks["summary"] + text_blocks["methods"])[:6]])
    add_textbox(slide, 0.7, 5.65, 11.3, 0.25, [{"text": f"Source type: {source['metadata']['type']} · Template asset: assets/template.pptx", "size": 11, "color": MID, "space_after": 0}], margin=0)

    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title(slide, "Method and Evidence")
    add_kicker(slide, "VALIDATION LOGIC")
    add_card(slide, 0.45, 1.55, 4.75, 4.3)
    add_textbox(slide, 0.7, 1.84, 4.25, 0.32, [{"text": "Method Summary", "size": 16, "bold": True, "space_after": 0}], margin=0)
    add_textbox(slide, 0.7, 2.2, 4.15, 3.3, [{"text": f"• {line}", "size": 13} for line in text_blocks["methods"][:5]])
    evidence_image = image_plan["method"]
    if evidence_image:
        add_picture_contain(slide, evidence_image["path"], 5.45, 1.55, 6.95, 3.15)
    else:
        add_placeholder(slide, 5.45, 1.55, 6.95, 3.15, "MEDIA / EVIDENCE PLACEHOLDER", "Insert test setup image, chart, or reference media here.", red=False)
    add_card(slide, 5.45, 4.95, 6.95, 0.9)
    method_caption = image_caption(evidence_image) if evidence_image else "Use this zone for setup notes, chart captions, or media commentary."
    add_textbox(slide, 5.75, 5.18, 6.3, 0.36, [{"text": method_caption, "size": 12, "space_after": 0}], margin=0)

    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title(slide, "Results and Interpretation")
    add_kicker(slide, "KEY OUTPUTS")
    add_card(slide, 0.45, 1.55, 12.0, 0.95)
    add_textbox(slide, 0.72, 1.82, 11.4, 0.36, [{"text": "The following slide consolidates extracted measured values and interprets them in professional client-facing language.", "size": 12, "color": DARK, "space_after": 0}], margin=0)
    add_card(slide, 0.45, 2.8, 4.15, 3.0)
    add_textbox(slide, 0.72, 3.08, 3.6, 2.45, [{"text": f"• {line}", "size": 13} for line in text_blocks["results"][:6]])
    add_card(slide, 4.82, 2.8, 4.25, 3.0)
    result_visual = image_plan["results"]
    if result_visual:
        add_picture_contain(slide, result_visual["path"], 5.02, 3.0, 3.85, 2.25)
        add_textbox(slide, 5.05, 5.33, 3.78, 0.18, [{"text": image_caption(result_visual), "size": 9, "color": MID, "align": PP_ALIGN.CENTER, "space_after": 0}], margin=0)
    else:
        add_placeholder(slide, 5.02, 3.0, 3.85, 2.25, "RESULT VISUAL", "Insert a chart, trend view, or comparison visual here.", red=False)
        add_textbox(slide, 5.05, 5.33, 3.78, 0.18, [{"text": "Use this zone for a result-supporting chart or comparison view.", "size": 9, "color": MID, "align": PP_ALIGN.CENTER, "space_after": 0}], margin=0)
    add_card(slide, 9.3, 2.8, 3.15, 3.0)
    metric_cards_vertical(slide, source.get("numeric_facts", []), left=9.55, top=3.0, width=2.65, height=2.56, max_cards=4)
    if richness >= 3 and not result_visual:
        add_textbox(slide, 9.58, 5.42, 2.58, 0.18, [{"text": "Prioritize measured values and validated evidence over decorative imagery.", "size": 8, "color": MID, "align": PP_ALIGN.CENTER, "space_after": 0}], margin=0)

    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title(slide, "Open Items and Placeholders")
    add_kicker(slide, "MEDIA COMPLETION")
    add_placeholder(slide, 0.45, 1.55, 5.85, 3.95, "VIDEO PLACEHOLDER", "Insert final product or test video here.")
    add_placeholder(slide, 6.55, 1.55, 5.85, 3.95, "IMAGE PLACEHOLDER", "Insert any missing thermal image, trend chart, or supporting figure here.", red=False)
    add_card(slide, 0.45, 5.72, 12.0, 0.45)
    add_textbox(slide, 0.72, 5.85, 11.3, 0.2, [{"text": "Use this slide whenever the source refers to video zones, missing images, or media that must be inserted later.", "size": 11, "color": MID, "space_after": 0}], margin=0)

    build_summary_slide(prs, "Conclusion", text_blocks["conclusion"], kicker="FINAL TAKEAWAY")
    build_thank_you_slide(prs)
    prs.save(str(output_path))


def build_market_research(prs: Presentation, title: str, source: dict, text_blocks: dict[str, list[str]], image_paths: list[Path], external_images: list[dict], output_path: Path, generated_visuals: list[Path], richness: int) -> None:
    image_plan = choose_market_images(image_paths, external_images, generated_visuals)
    extra_paths = remaining_catalog_paths(image_paths, external_images, generated_visuals, "market-research", [image_plan["landscape"], image_plan["implication"]])
    build_cover(prs, title, "Market research deck generated from the Sonceboz template")
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title(slide, "Executive Summary")
    add_kicker(slide, "MARKET OVERVIEW")
    add_card(slide, 0.45, 1.6, 12.0, 3.85)
    add_textbox(slide, 0.72, 1.9, 11.4, 3.15, [{"text": f"• {line}", "size": 15} for line in text_blocks["summary"][:6]])
    metric_cards_from_facts(slide, source.get("numeric_facts", []))
    if richness >= 4 and extra_paths and not filtered_metric_facts(source.get("numeric_facts", [])):
        add_textbox(slide, 0.72, 5.2, 2.2, 0.2, [{"text": "Visual Signal Strip", "size": 10, "bold": True, "color": MID, "space_after": 0}], margin=0)
        add_image_mosaic(slide, extra_paths, 0.72, 5.38, 11.1, 0.72, max_items=min(4, len(extra_paths)))

    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title(slide, "Landscape and Signals")
    add_kicker(slide, "MARKET LANDSCAPE")
    add_card(slide, 0.45, 1.55, 5.85, 4.1)
    add_card(slide, 6.55, 1.55, 5.85, 4.1)
    add_textbox(slide, 0.72, 1.84, 5.2, 3.3, [{"text": f"• {line}", "size": 14} for line in text_blocks["landscape"][:5]])
    market_image = image_plan["landscape"]
    if market_image:
        add_picture_contain(slide, market_image["path"], 6.78, 1.82, 5.35, 3.25)
        add_textbox(slide, 6.82, 5.08, 5.1, 0.2, [{"text": image_caption(market_image), "size": 9, "color": MID, "align": PP_ALIGN.CENTER, "space_after": 0}], margin=0)
    else:
        add_placeholder(slide, 6.78, 1.82, 5.35, 3.55, "VISUAL PLACEHOLDER", "Insert a competitor, product, or market visual.", red=False)

    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title(slide, "Implications and Recommendation")
    add_kicker(slide, "CLIENT TAKEAWAY")
    add_card(slide, 0.45, 1.55, 5.85, 4.1)
    add_card(slide, 6.55, 1.55, 5.85, 4.1)
    add_textbox(slide, 0.72, 1.84, 5.2, 3.3, [{"text": f"• {line}", "size": 14} for line in text_blocks["findings"][:5]])
    implication_image = image_plan["implication"]
    if richness >= 2 and implication_image:
        add_picture_contain(slide, implication_image["path"], 6.82, 1.84, 5.1, 3.05)
        add_textbox(slide, 6.9, 5.0, 4.95, 0.2, [{"text": image_caption(implication_image), "size": 9, "color": MID, "align": PP_ALIGN.CENTER, "space_after": 0}], margin=0)
    else:
        add_textbox(slide, 6.82, 1.84, 5.1, 3.3, [{"text": f"• {line}", "size": 14} for line in text_blocks["recommendation"][:5]])
    build_summary_slide(prs, "Recommendation", text_blocks["recommendation"], kicker="FINAL RECOMMENDATION")
    build_thank_you_slide(prs)
    prs.save(str(output_path))


def build_meeting_summary(prs: Presentation, title: str, source: dict, text_blocks: dict[str, list[str]], output_path: Path, generated_visuals: list[Path], external_images: list[dict], image_paths: list[Path]) -> None:
    image_plan = choose_meeting_images(image_paths, external_images, generated_visuals)
    build_cover(prs, title, "Meeting summary generated from the Sonceboz template")
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title(slide, "Summary")
    add_kicker(slide, "MEETING PURPOSE")
    add_card(slide, 0.45, 1.55, 12.0, 3.95)
    add_textbox(slide, 0.72, 1.9, 11.4, 3.15, [{"text": f"• {line}", "size": 15} for line in text_blocks["summary"][:6]])

    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title(slide, "Discussion Highlights")
    add_kicker(slide, "KEY POINTS")
    add_card(slide, 0.45, 1.55, 5.85, 4.1)
    add_card(slide, 6.55, 1.55, 5.85, 4.1)
    add_textbox(slide, 0.72, 1.84, 5.1, 3.3, [{"text": f"• {line}", "size": 14} for line in text_blocks["discussion"][:5]])
    meeting_image = image_plan["discussion"]
    if meeting_image:
        add_picture_contain(slide, meeting_image["path"], 6.82, 1.84, 5.1, 3.05)
        add_textbox(slide, 6.88, 4.98, 4.98, 0.22, [{"text": image_caption(meeting_image), "size": 9, "color": MID, "align": PP_ALIGN.CENTER, "space_after": 0}], margin=0)
    else:
        add_textbox(slide, 6.82, 1.84, 5.1, 3.3, [{"text": f"• {line}", "size": 14} for line in text_blocks["decisions"][:5]])

    build_summary_slide(prs, "Next Steps", text_blocks["next_steps"], kicker="ACTION SUMMARY")
    build_thank_you_slide(prs)
    prs.save(str(output_path))


def inspect_ppt(path: Path) -> None:
    prs = Presentation(str(path))
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = " ".join(part.strip() for part in shape.text.splitlines() if part.strip())
                if text:
                    texts.append(text)
        print(f"SLIDE {idx}: {' | '.join(texts[:12])}")


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--source", help="normalized JSON from extract_source.py")
    parser.add_argument("--purpose", choices=["test-report", "market-research", "meeting-summary"])
    parser.add_argument("--title", default="")
    parser.add_argument("--output", help="output pptx path")
    parser.add_argument("--inspect", help="inspect an existing pptx and print text")
    parser.add_argument("--auto-visuals", action="store_true", help="generate SVG/PNG supporting visuals and use them when suitable")
    parser.add_argument("--illustration-richness", type=int, default=0, help="0-5 extra illustration richness: 0 source/svg only, 5 web+AI rich")
    args = parser.parse_args()

    if args.inspect:
        inspect_ppt(Path(args.inspect).resolve())
        return

    if not args.source or not args.output or not args.purpose:
        raise SystemExit("--source, --purpose, and --output are required unless --inspect is used.")

    source = json.loads(Path(args.source).read_text(encoding="utf-8"))
    output_path = Path(args.output).resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    text_blocks = choose_lines_for_purpose(source.get("text", ""), args.purpose)
    image_paths = [Path(item) for item in source.get("image_paths", []) if Path(item).exists()]
    generated_visuals = generate_visual_assets(source, args.purpose, output_path) if args.auto_visuals else []
    richness = max(0, min(5, args.illustration_richness))
    external_images = generate_external_image_records(source, args.purpose, output_path, richness)
    presentation = Presentation(str(TEMPLATE_PATH))
    remove_all_slides(presentation)
    title = polish_title(args.title, args.purpose)
    source["title"] = title

    if args.purpose == "test-report":
        build_test_report(presentation, title, source, text_blocks, image_paths, external_images, output_path, generated_visuals, richness)
    elif args.purpose == "market-research":
        build_market_research(presentation, title, source, text_blocks, image_paths, external_images, output_path, generated_visuals, richness)
    else:
        build_meeting_summary(presentation, title, source, text_blocks, output_path, generated_visuals, external_images, image_paths)

    print(output_path)


if __name__ == "__main__":
    main()
