from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation


SKILL_DIR = Path(__file__).resolve().parents[1]
TEMPLATE_PATH = SKILL_DIR / "assets" / "template.pptx"


def main() -> None:
    presentation = Presentation(str(TEMPLATE_PATH))
    report: dict[str, object] = {
        "template": str(TEMPLATE_PATH),
        "slide_size": {
            "width": presentation.slide_width,
            "height": presentation.slide_height,
        },
        "layouts": [],
    }
    for index, layout in enumerate(presentation.slide_layouts):
        placeholders = []
        for shape in layout.shapes:
            item = {
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
            }
            if shape.is_placeholder:
                item["placeholder_idx"] = shape.placeholder_format.idx
                item["placeholder_type"] = str(shape.placeholder_format.type)
            placeholders.append(item)
        report["layouts"].append(
            {
                "index": index,
                "name": layout.name,
                "placeholders": placeholders,
            }
        )
    print(json.dumps(report, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
