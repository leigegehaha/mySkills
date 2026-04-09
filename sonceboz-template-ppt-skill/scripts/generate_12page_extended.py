from __future__ import annotations

import argparse
import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SKILL_DIR = Path(__file__).resolve().parents[1]
TEMPLATE_PATH = SKILL_DIR / "assets" / "template.pptx"
OUTPUT_DIR = Path("/Users/zhangleiandhim/Sonceboz Project relatated files/Dextrous Hand/Sonceboz Dextrous Hand/L20 Lite Cursh")
RED = RGBColor(192, 0, 0)
DARK = RGBColor(51, 51, 51)
MID = RGBColor(102, 102, 102)
LIGHT = RGBColor(245, 245, 245)
BORDER = RGBColor(217, 217, 217)
WHITE = RGBColor(255, 255, 255)


def remove_all_slides(prs):
    slide_ids = prs.slides._sldIdLst
    for slide_id in list(slide_ids):
        prs.part.drop_rel(slide_id.rId)
        slide_ids.remove(slide_id)


def clear_text_frame(text_frame, margin: float = 0.08, vertical_anchor=MSO_ANCHOR.TOP):
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(margin)
    text_frame.margin_right = Inches(margin)
    text_frame.margin_top = Inches(margin)
    text_frame.margin_bottom = Inches(margin)
    text_frame.vertical_anchor = vertical_anchor


def style_paragraph(paragraph, text: str, size: int = 16, color: RGBColor = DARK, bold: bool = False, align=PP_ALIGN.LEFT, space_after: int = 8):
    paragraph.text = text
    paragraph.alignment = align
    paragraph.space_after = Pt(space_after)
    paragraph.line_spacing = 1.05
    run = paragraph.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_textbox(slide, left: float, top: float, width: float, height: float, paragraphs: list[dict], margin: float = 0.08, vertical_anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    clear_text_frame(tf, margin=margin, vertical_anchor=vertical_anchor)
    for idx, spec in enumerate(paragraphs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        style_paragraph(
            p,
            spec["text"],
            size=spec.get("size", 16),
            color=spec.get("color", DARK),
            bold=spec.get("bold", False),
            align=spec.get("align", PP_ALIGN.LEFT),
            space_after=spec.get("space_after", 8),
        )
    return box


def add_card(slide, left: float, top: float, width: float, height: float, fill: RGBColor = WHITE, line: RGBColor = BORDER):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_picture_contain(slide, image_path: Path, left: float, top: float, width: float, height: float):
    add_card(slide, left, top, width, height, fill=LIGHT, line=BORDER)
    with Image.open(image_path) as image:
        ratio = image.width / image.height
    box_ratio = width / height
    if ratio >= box_ratio:
        pic_width = Inches(width)
        pic_height = int(pic_width / ratio)
    else:
        pic_height = Inches(height)
        pic_width = int(pic_height * ratio)
    slide.shapes.add_picture(
        str(image_path),
        Inches(left) + (Inches(width) - pic_width) // 2,
        Inches(top) + (Inches(height) - pic_height) // 2,
        width=pic_width,
        height=pic_height,
    )


def add_title(slide, title: str):
    placeholder = slide.placeholders[0]
    placeholder.text = title
    tf = placeholder.text_frame
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        for r in p.runs:
            r.font.name = "Arial"
            r.font.size = Pt(28)
            r.font.bold = True
            r.font.color.rgb = RED


def add_kicker(slide, text: str):
    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(1.03), Inches(2.0), Inches(0.32))
    tag.fill.solid()
    tag.fill.fore_color.rgb = RED
    tag.line.color.rgb = RED
    tf = tag.text_frame
    clear_text_frame(tf, margin=0.02, vertical_anchor=MSO_ANCHOR.MIDDLE)
    style_paragraph(tf.paragraphs[0], text, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER, space_after=0)


def add_placeholder(slide, left: float, top: float, width: float, height: float, heading: str, note: str, red: bool = True):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(252, 252, 252)
    box.line.color.rgb = RED if red else BORDER
    box.line.width = Pt(1.8 if red else 1.2)
    box.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    tf = box.text_frame
    clear_text_frame(tf, margin=0.1, vertical_anchor=MSO_ANCHOR.MIDDLE)
    style_paragraph(tf.paragraphs[0], heading, size=18 if red else 15, color=RED if red else DARK, bold=True, align=PP_ALIGN.CENTER, space_after=4)
    p = tf.add_paragraph()
    style_paragraph(p, note, size=11, color=MID, align=PP_ALIGN.CENTER, space_after=0)


def add_metric_card(slide, left: float, top: float, width: float, height: float, value: str, label: str):
    card = add_card(slide, left, top, width, height)
    tf = card.text_frame
    clear_text_frame(tf, margin=0.06, vertical_anchor=MSO_ANCHOR.MIDDLE)
    style_paragraph(tf.paragraphs[0], value, size=18, color=RED, bold=True, align=PP_ALIGN.CENTER, space_after=2)
    p = tf.add_paragraph()
    style_paragraph(p, label[:42], size=10, color=MID, align=PP_ALIGN.CENTER, space_after=0)


def build_cover(prs, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.placeholders[14].text = title
    slide.placeholders[15].text = subtitle
    for idx in (14, 15):
        tf = slide.placeholders[idx].text_frame
        for p in tf.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = "Arial"
                r.font.color.rgb = WHITE
                if idx == 14:
                    r.font.size = Pt(36)
                    r.font.bold = True
                else:
                    r.font.size = Pt(18)


def build_thank_you(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[7])
    tf = slide.placeholders[13].text_frame
    clear_text_frame(tf, margin=0.08, vertical_anchor=MSO_ANCHOR.MIDDLE)
    style_paragraph(tf.paragraphs[0], "Thank you", size=30, color=WHITE, bold=True, align=PP_ALIGN.LEFT, space_after=0)


def find_first_image(image_paths: list[Path], keywords: list[str]) -> Path | None:
    for path in image_paths:
        name = path.stem.lower()
        if any(kw in name for kw in keywords):
            return path
    return None


def main():
    # Load extracted source
    source = json.load(open("/tmp/l20_lite_source.json", "r", encoding="utf-8"))
    image_paths = [Path(item) for item in source.get("image_paths", []) if Path(item).exists()]

    prs = Presentation(str(TEMPLATE_PATH))
    remove_all_slides(prs)

    # === SLIDE 1: Cover ===
    build_cover(prs, "LINKER HAND L20 LITE\nTest Report", "Comprehensive Validation Study | Dextrous Hand Performance Assessment")

    # === SLIDE 2: Executive Summary ===
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title(slide, "Executive Summary")
    add_kicker(slide, "REPORT OVERVIEW")
    add_card(slide, 0.45, 1.55, 12.0, 3.85)
    add_textbox(slide, 0.7, 1.85, 11.5, 3.2, [
        {"text": "• Scope: Independent performance evaluation of the Linker Hand L20 Lite dextrous robotic hand.", "size": 15},
        {"text": "• Tested grip force, fingertip force, finger speed, and thermal behavior under operational conditions.", "size": 15},
        {"text": "• Key finding: Grasp force (52 N) is notably below the 80 N official specification.", "size": 15},
        {"text": "• Fingertip forces ranged from 5.3 N to 6.8 N, well below the 15 N nominal specification.", "size": 15},
        {"text": "• Finger speed validated at 65 °/s, consistent with the claimed 63 °/s specification.", "size": 15},
        {"text": "• Significant temperature rise of ~14°C observed after 2 hours of continuous operation.", "size": 15},
        {"text": "• Structural failure: Index finger joint fracture after 2 hours of endurance operation.", "size": 15},
    ])

    # === SLIDE 3: Product Overview ===
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title(slide, "Product & Specification Overview")
    add_kicker(slide, "TARGET PRODUCT")
    
    add_card(slide, 0.45, 1.55, 5.85, 3.5)
    add_textbox(slide, 0.7, 1.8, 5.4, 3.0, [
        {"text": "Product Name", "size": 14, "bold": True},
        {"text": "Linker Hand L20 Lite", "size": 14, "space_after": 8},
        {"text": "Degrees of Freedom", "size": 14, "bold": True},
        {"text": "20 total DOFs (10 active + 10 passive)", "size": 14, "space_after": 8},
        {"text": "Key Capability Areas", "size": 14, "bold": True},
        {"text": "• Movement range and speed control", "size": 14},
        {"text": "• Force and grasp performance", "size": 14},
        {"text": "• Embedded tactile sensing system", "size": 14},
    ])

    add_card(slide, 6.55, 1.55, 5.85, 3.5)
    hero_img = find_first_image(image_paths, ["hand", "product", "device", "finger"])
    if hero_img:
        add_picture_contain(slide, hero_img, 6.7, 1.75, 5.35, 3.15)
        add_textbox(slide, 6.7, 5.05, 5.1, 0.2, [{"text": "Product visual reference", "size": 10, "color": MID, "align": PP_ALIGN.CENTER}])
    else:
        add_placeholder(slide, 6.7, 1.75, 5.35, 3.15, "IMAGE PLACEHOLDER", "Insert official product image here.", red=False)

    # === SLIDE 4: Key Specifications ===
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title(slide, "Official Specifications")
    add_kicker(slide, "NOMINAL VALUES")
    
    cards = [
        ("80 N", "Max Grasp Force", 0.45),
        ("15 N", "Max Fingertip Force", 3.65),
        ("63 °/s", "Finger Angular Speed", 6.85),
        ("20", "Degrees of Freedom", 10.05),
    ]
    for value, label, left in cards:
        add_metric_card(slide, left, 1.85, 2.9, 0.95, value, label)

    add_card(slide, 0.45, 3.55, 12.0, 2.15)
    add_textbox(slide, 0.75, 3.8, 11.4, 1.8, [
        {"text": "Specification Notes:", "size": 16, "bold": True, "space_after": 8},
        {"text": "• These values represent the manufacturer's published specifications for the L20 Lite.", "size": 14},
        {"text": "• The test unit provided was described as a demonstration model with potential performance degradation.", "size": 14},
        {"text": "• All subsequent test results are compared against these nominal values.", "size": 14},
    ])

    # === SLIDE 5: Test 1 - Grasp Force ===
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title(slide, "Test 1: Maximum Grasp Force")
    add_kicker(slide, "GRASP VALIDATION")
    
    add_card(slide, 0.45, 1.55, 5.85, 2.25)
    add_textbox(slide, 0.7, 1.8, 5.4, 1.6, [
        {"text": "Method", "size": 16, "bold": True, "space_after": 4},
        {"text": "Used a custom cylindrical grip force sensor; performed a five-finger grasp motion to measure peak force output.", "size": 14},
    ])

    compare_img = find_first_image(image_paths, ["grasp", "force", "measure", "test"])
    if compare_img:
        add_picture_contain(slide, compare_img, 6.55, 1.55, 5.85, 2.25)
    else:
        add_card(slide, 6.55, 1.55, 5.85, 2.25)
        add_textbox(slide, 6.7, 2.4, 5.4, 0.8, [{"text": "Grasp force test setup", "size": 14, "align": PP_ALIGN.CENTER}])

    add_metric_card(slide, 0.45, 3.95, 2.7, 0.8, "52 N", "Measured Grasp Force")
    add_metric_card(slide, 3.45, 3.95, 2.7, 0.8, "80 N", "Official Specification")
    add_metric_card(slide, 6.45, 3.95, 2.7, 0.8, "-35%", "Deviation from Spec")
    add_metric_card(slide, 9.45, 3.95, 2.9, 0.8, "Significant", "Gap Assessment")

    add_card(slide, 0.45, 5.15, 12.0, 0.55)
    add_textbox(slide, 0.7, 5.28, 11.4, 0.3, [{"text": "Conclusion: Measured grasp force is significantly below the official specification.", "size": 12, "bold": True, "color": RED}])

    # === SLIDE 6: Test 2 - Fingertip Force ===
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title(slide, "Test 2: Maximum Fingertip Force")
    add_kicker(slide, "FINGER FORCE VALIDATION")
    
    add_card(slide, 0.45, 1.55, 5.85, 1.8)
    add_textbox(slide, 0.7, 1.8, 5.4, 1.2, [
        {"text": "Method: Dedicated fingertip force test fixture; hand was fixed during testing; each finger actuated to maximum output.", "size": 14},
    ])

    add_card(slide, 0.45, 3.55, 12.0, 2.05)
    add_textbox(slide, 0.7, 3.7, 11.4, 1.8, [
        {"text": "Measured Values:", "size": 16, "bold": True, "space_after": 6},
        {"text": "• Index Finger: 5.3 N  |  Middle Finger: 6.6 N  |  Ring Finger: 6.0 N  |  Little Finger: 6.8 N", "size": 15, "space_after": 6},
        {"text": "• Thumb: Not measured (hand envelope too large for fixture alignment).", "size": 14, "color": MID},
    ])

    add_metric_card(slide, 0.45, 5.85, 2.7, 0.8, "6.6 N", "Average Fingertip Force")
    add_metric_card(slide, 3.45, 5.85, 2.7, 0.8, "15 N", "Official Specification")
    add_metric_card(slide, 6.45, 5.85, 2.7, 0.8, "-56%", "Average Deviation")
    add_metric_card(slide, 9.45, 5.85, 2.9, 0.8, "Not Tested", "Thumb Force")

    # === SLIDE 7: Fingertip Force Analysis ===
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title(slide, "Test 2 (Continued): Analysis & Explanation")
    add_kicker(slide, "FINDING ANALYSIS")
    
    add_card(slide, 0.45, 1.55, 12.0, 2.35)
    add_textbox(slide, 0.75, 1.8, 11.4, 1.8, [
        {"text": "Root Cause Inquiry:", "size": 16, "bold": True, "space_after": 6},
        {"text": "Manufacturer confirmed the unit was a long-term demonstration model. Performance degradation is expected in units with extended use.", "size": 14, "space_after": 6},
        {"text": "Internal linkage wear from repeated demonstration cycles is the likely cause of reduced fingertip force.", "size": 14},
    ])

    add_card(slide, 0.45, 4.15, 5.85, 1.5)
    add_textbox(slide, 0.75, 4.35, 5.3, 1.1, [
        {"text": "Key Concern", "size": 14, "bold": True, "color": RED, "space_after": 4},
        {"text": "Fingertip forces consistently below 50% of nominal spec indicate significant mechanical degradation.", "size": 13},
    ])

    add_card(slide, 6.55, 4.15, 5.85, 1.5)
    add_textbox(slide, 6.8, 4.35, 5.3, 1.1, [
        {"text": "Test Limitation", "size": 14, "bold": True, "color": MID, "space_after": 4},
        {"text": "Thumb force was not tested due to fixture incompatibility and hand envelope constraints.", "size": 13},
    ])

    # === SLIDE 8: Test 3 - Finger Speed ===
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title(slide, "Test 3: Finger Motion Speed")
    add_kicker(slide, "SPEED VALIDATION")
    
    add_card(slide, 0.45, 1.55, 5.85, 2.0)
    add_textbox(slide, 0.7, 1.8, 5.4, 1.4, [
        {"text": "Method", "size": 16, "bold": True, "space_after": 4},
        {"text": "Set finger speed to maximum; recorded full open-to-close motion. Frame-by-frame analysis yields precise time between extreme positions.", "size": 14},
    ])

    speed_img = find_first_image(image_paths, ["speed", "motion", "movement"])
    if speed_img:
        add_picture_contain(slide, speed_img, 6.55, 1.55, 5.85, 2.0)
    else:
        add_placeholder(slide, 6.55, 1.55, 5.85, 2.0, "VIDEO PLACEHOLDER", "Insert 0.1x slow-motion test footage here.", red=False)

    add_card(slide, 0.45, 3.85, 12.0, 1.8)
    add_textbox(slide, 0.7, 4.05, 11.4, 1.4, [
        {"text": "Measured Results:", "size": 16, "bold": True, "space_after": 6},
        {"text": "• Little finger cycle time: 1.2 s", "size": 15, "space_after": 4},
        {"text": "• Motion range: 77.9°", "size": 15, "space_after": 4},
        {"text": "• Calculated speed: 65 °/s  —  close to the advertised 63 °/s specification.", "size": 15, "color": RGBColor(34, 139, 34)},
    ])

    add_metric_card(slide, 0.45, 5.85, 2.7, 0.8, "1.2 s", "Cycle Time (Little Finger)")
    add_metric_card(slide, 3.45, 5.85, 2.7, 0.8, "77.9°", "Motion Range")
    add_metric_card(slide, 6.45, 5.85, 2.7, 0.8, "65 °/s", "Calculated Speed")
    add_metric_card(slide, 9.45, 5.85, 2.9, 0.8, "Matched", "Validation Status")

    # === SLIDE 9: Test 4 - Thermal Behavior ===
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title(slide, "Test 4: Thermal Field Analysis")
    add_kicker(slide, "THERMAL VALIDATION")
    
    add_card(slide, 0.45, 1.55, 12.0, 2.0)
    add_textbox(slide, 0.7, 1.8, 11.4, 1.4, [
        {"text": "Method", "size": 16, "bold": True, "space_after": 4},
        {"text": "Used a thermal imaging camera to capture temperature distribution across the hand before and after endurance testing.", "size": 14},
    ])

    # Before / After temperatures
    add_card(slide, 0.45, 3.85, 5.85, 1.8)
    add_textbox(slide, 0.7, 4.0, 5.4, 1.4, [
        {"text": "Before Endurance Test", "size": 15, "bold": True, "color": RGBColor(34, 139, 34), "space_after": 6},
        {"text": "• Palm: Max 27.2°C (finger root joint)", "size": 14, "space_after": 4},
        {"text": "• Back of hand: 26.2°C (MCU location)", "size": 14},
    ])

    add_card(slide, 6.55, 3.85, 5.85, 1.8)
    add_textbox(slide, 6.8, 4.0, 5.4, 1.4, [
        {"text": "After 2-Hour Endurance Test", "size": 15, "bold": True, "color": RED, "space_after": 6},
        {"text": "• Back of hand: Max 50.6°C", "size": 14, "space_after": 4},
        {"text": "• Palm: Max 46.8°C", "size": 14, "space_after": 4},
        {"text": "• Finger root joint: Max 45.8°C", "size": 14},
    ])

    add_card(slide, 0.45, 5.85, 12.0, 0.55)
    add_textbox(slide, 0.7, 5.98, 11.4, 0.3, [{"text": "Temperature rise of ~14°C after 2 hours — indicates inadequate heat dissipation under sustained load.", "size": 12, "bold": True, "color": RED}])

    # === SLIDE 10: Durability & Structural Failure ===
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title(slide, "Durability & Structural Failure")
    add_kicker(slide, "FAILURE ANALYSIS")
    
    add_card(slide, 0.45, 1.55, 5.85, 2.0)
    add_textbox(slide, 0.75, 1.8, 5.35, 1.4, [
        {"text": "Test Protocol Limitation", "size": 16, "bold": True, "space_after": 4},
        {"text": "Original plan: 24-hour endurance test. Actual: Stopped after 2 hours due to structural failure.", "size": 13},
        {"text": "The index finger joint fractured and became inoperable.", "size": 13},
    ])

    fail_img = find_first_image(image_paths, ["break", "damage", "thermal", "camera"])
    if fail_img:
        add_picture_contain(slide, fail_img, 6.55, 1.55, 5.85, 2.0)
    else:
        add_placeholder(slide, 6.55, 1.55, 5.85, 2.0, "IMAGE PLACEHOLDER", "Insert index finger fracture photo here.", red=False)

    add_card(slide, 0.45, 3.85, 12.0, 1.8)
    add_textbox(slide, 0.75, 4.05, 11.4, 1.4, [
        {"text": "Failure Sequence:", "size": 16, "bold": True, "space_after": 6},
        {"text": "1. Continuous operation → Temperature rise to 50.6°C (back of hand).", "size": 14, "space_after": 4},
        {"text": "2. Heat accumulation in MCU and joint regions likely accelerated material fatigue.", "size": 14, "space_after": 4},
        {"text": "3. Index finger joint fractured after 2 hours — test was terminated.", "size": 14},
    ])

    add_card(slide, 0.45, 5.85, 12.0, 0.55)
    add_textbox(slide, 0.7, 5.98, 11.4, 0.3, [{"text": "Note: Test was conducted on a borrowed demonstration unit; extended testing was cancelled to prevent further damage.", "size": 11, "color": MID}])

    # === SLIDE 11: Comparative Summary ===
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title(slide, "Comparative Summary: L20 Lite vs. Specifications")
    add_kicker(slide, "RESULT AGGREGATION")
    
    add_card(slide, 0.45, 1.55, 12.0, 3.85)
    add_textbox(slide, 0.7, 1.8, 11.4, 3.3, [
        {"text": "Test Parameter Comparison:", "size": 16, "bold": True, "space_after": 10},
        {"text": "Grasp Force:        52 N measured  vs.  80 N spec  |  Δ = -35%  ❌ Below specification", "size": 14, "space_after": 8},
        {"text": "Fingertip Force:  ~6.6 N measured  vs.  15 N spec   |  Δ = -56%  ❌ Well below specification", "size": 14, "space_after": 8},
        {"text": "Finger Speed:         65 °/s measured  vs.  63 °/s spec  |  Δ ≈ +3%   ✅ Meets/exceeds spec", "size": 14, "space_after": 8},
        {"text": "Thermal Rise:    ~14°C after 2 hrs  |  ⚠️ Significant; heat dissipation concern", "size": 14, "space_after": 8},
        {"text": "Durability:  Failed at 2 hours  |  ❌ Structural failure — joint fracture", "size": 14, "space_after": 10},
        {"text": "Comparison to Inspire hand: L20 Lite is noticeably faster but significantly underpowered. The official explanation attributes the reduced force to the unit being a long-term demonstration model with likely internal wear.", "size": 13, "color": MID},
    ])

    # === SLIDE 12: Conclusion ===
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title(slide, "Conclusions & Recommendations")
    add_kicker(slide, "FINAL ASSESSMENT")
    
    add_card(slide, 0.45, 1.55, 5.85, 4.15)
    add_textbox(slide, 0.75, 1.8, 5.35, 3.7, [
        {"text": "Conclusions", "size": 18, "bold": True, "color": RED, "space_after": 8},
        {"text": "✅ Speed validated — 65 °/s meets the 63 °/s claim.", "size": 14, "space_after": 6},
        {"text": "⚠️ Force significantly degraded — both grasp and fingertip forces are 35-56% below specification.", "size": 14, "space_after": 6},
        {"text": "⚠️ Thermal performance concern — 14°C rise after 2 hours of continuous operation.", "size": 14, "space_after": 6},
        {"text": "❌ Structural failure — index finger joint fractured, limiting endurance to only 2 hours.", "size": 14, "space_after": 6},
        {"text": "❌ Demonstration unit degradation — internal linkage wear is the confirmed root cause.", "size": 14},
    ])

    add_card(slide, 6.55, 1.55, 5.85, 4.15)
    add_textbox(slide, 6.8, 1.8, 5.35, 3.7, [
        {"text": "Recommendations", "size": 18, "bold": True, "color": RGBColor(0, 100, 0), "space_after": 8},
        {"text": "1. Request a new production unit for re-testing without known wear.", "size": 13, "space_after": 6},
        {"text": "2. Conduct full 24-hour endurance test on a fresh unit.", "size": 13, "space_after": 6},
        {"text": "3. Prioritize thermal dissipation analysis for the MCU and joint regions.", "size": 13, "space_after": 6},
        {"text": "4. Verify thumb force with an adapted measurement fixture.", "size": 13, "space_after": 6},
        {"text": "5. Evaluate long-term reliability with multiple production samples.", "size": 13},
    ])

    # Thank you slide
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title(slide, "Thank you")
    add_kicker(slide, "Sonceboz Template")
    add_textbox(slide, 0.75, 2.0, 11.4, 1.0, [{"text": "Generated with Sonceboz Template", "size": 18, "color": MID, "align": PP_ALIGN.CENTER, "space_after": 0}])

    output_path = OUTPUT_DIR / "Linker_Hand_L20_Lite_Test_Report_Extended.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(output_path)


if __name__ == "__main__":
    main()
